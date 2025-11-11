# -*- coding: utf-8 -*-
"""
To‘liq ishlaydigan Lotin ↔ Krill bot (docx, xlsx, pptx formatlarda)
Talablar (requirements):
    python-docx, openpyxl, python-pptx, python-telegram-bot
"""

import logging
import os
import re
import tempfile
import shutil
from telegram import Update, ReplyKeyboardMarkup
from telegram.ext import (
    ApplicationBuilder,
    CommandHandler,
    MessageHandler,
    filters,
    ContextTypes,
)
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

# --- CONFIG ---
TOKEN = "8451093366:AAFr9rkhFJYV060TPcjB069iKDd6sObezXQ"

# --- Logging ---
logging.basicConfig(format="%(asctime)s - %(levelname)s - %(message)s", level=logging.INFO)
logger = logging.getLogger(__name__)

# --- Normalize apostrophes ---
def normalize_apostrophes(text: str) -> str:
    """Turli apostrof belgilarini yagona `'` ga almashtiradi."""
    if not text:
        return text
    for a in ["‘", "’", "ʼ", "ʻ", "`", "ʹ", "´", "′"]:
        text = text.replace(a, "'")
    return text

# --- Force e rule (Krill → Lotin) ---
def force_e_rule(text: str) -> str:
    """
    Krill → Lotin transliteratsiyasida 'е' harfi uchun to‘g‘ri qoida:
    - So‘z boshida yoki unli, ъ, ь dan keyin kelsa → 'ye'
    - Aks holda → 'e'
    - 'Э' harfi esa doim 'e'
    """
    text = text.replace('Э', 'E').replace('э', 'e')

    # So‘z boshida yoki unli, ъ, ь dan keyin 'е' → 'ye'
    text = re.sub(r'\bЕ', 'Ye', text)
    text = re.sub(r'\bе', 'ye', text)
    text = re.sub(
        r'([АаЕеЁёИиОоУуЮюЯяЪъЬь])Е',
        lambda m: m.group(1) + 'Ye',
        text
    )
    text = re.sub(
        r'([АаЕеЁёИиОоУуЮюЯяЪъЬь])е',
        lambda m: m.group(1) + 'ye',
        text
    )

    # Qolgan 'е' → 'e'
    text = text.replace('Е', 'E').replace('е', 'e')
    return text

# --- Apostrophe rules for oʻ & gʻ ---
def apply_apostrophe_rules(s: str) -> str:
    s = normalize_apostrophes(s)
    s = re.sub(r"[Oo]'(?=[^a-zA-Z]|$)", lambda m: "ў" if m.group(0)[0].islower() else "Ў", s)
    s = re.sub(r"[Gg]'(?=[^a-zA-Z]|$)", lambda m: "ғ" if m.group(0)[0].islower() else "Ғ", s)
    return s

# --- Transliteration maps ---
LATIN_TO_CYR = {
    "sh": "ш", "ch": "ч", "ts": "ц",
    "a":"а","b":"б","d":"д","e":"е","f":"ф","g":"г","h":"ҳ","i":"и",
    "j":"ж","k":"к","l":"л","m":"м","n":"н","o":"о","p":"п","q":"қ","r":"р",
    "s":"с","t":"т","u":"у","v":"в","x":"х","y":"й","z":"з",
    "yo": "ё", "ya": "я", "yu": "ю", "ye": "е",
}
CYR_TO_LAT = {v:k for k,v in LATIN_TO_CYR.items()}

# Krill → Lotin o‘zgarishlari
CYR_TO_LAT['ц'] = 's'
CYR_TO_LAT['Ц'] = 'S'
CYR_TO_LAT['ъ'] = '’'
CYR_TO_LAT['Ъ'] = '’'

def expand_case(d):
    out = {}
    for k,v in d.items():
        out[k] = v
        out[k.upper()] = v.upper()
        if len(k) > 1:
            out[k.capitalize()] = v.capitalize()
    return out

LATIN_TO_CYR = expand_case(LATIN_TO_CYR)
CYR_TO_LAT = expand_case(CYR_TO_LAT)

LAT_TO_CYR_RE = re.compile("|".join(sorted(map(re.escape, LATIN_TO_CYR.keys()), key=len, reverse=True)))
CYR_TO_LAT_RE = re.compile("|".join(sorted(map(re.escape, CYR_TO_LAT.keys()), key=len, reverse=True)))

# --- Replace match case (faqat bosh harf capitalize, CH/SH bo'lmasin) ---
def replace_match_case(m, mapping, start_of_sentence=False):
    t = m.group(0)
    if start_of_sentence and t.lower() in ["ch", "sh", "ts", "yo", "ya", "yu", "ye"]:
        return t.capitalize()  # Ch, Sh
    if t[0].isupper():
        return mapping.get(t.lower(), t).capitalize()
    return mapping.get(t, t)

# --- Transliteration function with sentence-aware Ch/Sh ---
def transliterate_text(s: str, to_cyr=True) -> str:
    if not s:
        return s
    s = normalize_apostrophes(s)

    if to_cyr:
        s = apply_apostrophe_rules(s)
        s = re.sub(r'\b[eE]', lambda m: 'э' if m.group(0).islower() else 'Э', s)
        s = re.sub(r'\b[hH]', lambda m: 'ҳ' if m.group(0).islower() else 'Ҳ', s)
        s = LAT_TO_CYR_RE.sub(lambda m: replace_match_case(m, LATIN_TO_CYR), s)
    else:
        s = force_e_rule(s)

        # Gaplarni ajratib, gap boshidagi so‘zlarni aniqlaymiz
        parts = re.split(r'([.!?]\s*)', s)
        new_parts = []
        start_sentence = True
        for p in parts:
            if re.match(r'[.!?]\s*', p):
                new_parts.append(p)
                start_sentence = True
            else:
                # har so‘zni alohida transliteratsiya qilamiz
                words = re.findall(r'\w+|\W+', p)
                new_words = []
                for w in words:
                    if re.match(r'\w+', w):
                        w_conv = CYR_TO_LAT_RE.sub(lambda m: replace_match_case(m, CYR_TO_LAT, start_sentence), w)
                        w_conv = w_conv.replace("ў", "o‘").replace("Ў", "O‘")
                        w_conv = w_conv.replace("ғ", "g‘").replace("Ғ", "G‘")
                        w_conv = w_conv.replace("ҳ", "h").replace("Ҳ", "H")
                        new_words.append(w_conv)
                        start_sentence = False
                    else:
                        new_words.append(w)
                new_parts.append(''.join(new_words))
        s = ''.join(new_parts)

    return s

# --- DOCX conversion ---
def convert_docx_preserve_format(input_path, output_path, to_cyr=True):
    doc = Document(input_path)
    for para in doc.paragraphs:
        for run in para.runs:
            if run.text:
                run.text = transliterate_text(run.text, to_cyr)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        if run.text:
                            run.text = transliterate_text(run.text, to_cyr)
    doc.save(output_path)
    return output_path

# --- EXCEL conversion ---
def convert_xlsx_preserve_format(input_path, output_path, to_cyr=True):
    wb = load_workbook(input_path)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                val = cell.value
                if isinstance(val, str):
                    cell.value = transliterate_text(val, to_cyr)
    wb.save(output_path)
    return output_path

# --- POWERPOINT conversion ---
def convert_pptx_preserve_format(input_path, output_path, to_cyr=True):
    prs = Presentation(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            run.text = transliterate_text(run.text, to_cyr)
    prs.save(output_path)
    return output_path

# --- Telegram handlers ---
START_TEXT = "Assalomu alaykum!\nMen Lotin ↔ Krill botiman.\nFayl yuboring (.docx, .xlsx, .pptx) yoki matn yozing."

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    kb = [["Lotindan Krillga", "Krilldan Lotinga"]]
    await update.message.reply_text(START_TEXT, reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True))

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    mode = context.user_data.get("mode")
    if not mode:
        await update.message.reply_text("Avval yo‘nalishni tanlang: Lotindan Krillga yoki Krilldan Lotinga.")
        return
    doc = update.message.document
    fname = doc.file_name.lower()
    file = await context.bot.get_file(doc.file_id)
    tmp_dir = tempfile.mkdtemp(prefix="tg_doc_")
    in_path = os.path.join(tmp_dir, fname)
    out_path = os.path.join(tmp_dir, "converted_" + fname)
    await file.download_to_drive(in_path)

    to_cyr = (mode == "to_cyr")
    try:
        if fname.endswith(".docx"):
            convert_docx_preserve_format(in_path, out_path, to_cyr)
        elif fname.endswith(".xlsx"):
            convert_xlsx_preserve_format(in_path, out_path, to_cyr)
        elif fname.endswith(".pptx"):
            convert_pptx_preserve_format(in_path, out_path, to_cyr)
        else:
            await update.message.reply_text("❗ Faqat .docx, .xlsx, .pptx fayllarni qabul qilaman.")
            return
        with open(out_path, "rb") as f:
            await update.message.reply_document(f, filename=os.path.basename(out_path))
    except Exception as e:
        logger.exception("Xatolik:")
        await update.message.reply_text("❌ Faylni o‘zgartirishda xatolik yuz berdi.")
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

async def text_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    text = update.message.text or ""
    if text == "Lotindan Krillga":
        context.user_data["mode"] = "to_cyr"
        await update.message.reply_text("✅ Lotindan Krillga rejimi tanlandi.")
        return
    if text == "Krilldan Lotinga":
        context.user_data["mode"] = "to_lat"
        await update.message.reply_text("✅ Krilldan Lotinga rejimi tanlandi.")
        return
    mode = context.user_data.get("mode")
    if not mode:
        await update.message.reply_text("Avval yo‘nalishni tanlang.")
        return
    to_cyr = (mode == "to_cyr")
    result = transliterate_text(text, to_cyr)
    await update.message.reply_text(result)

async def error_handler(update: object, context: ContextTypes.DEFAULT_TYPE) -> None:
    logger.exception("Unhandled exception:")
    try:
        if hasattr(update, "message") and update.message:
            await update.message.reply_text("Kutilmagan xatolik yuz berdi.")
    except Exception:
        pass

def main():
    app = ApplicationBuilder().token(TOKEN).build()
    app.add_handler(CommandHandler("start", start))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, text_message))
    app.add_error_handler(error_handler)
    logger.info("Bot ishga tushdi...")
    app.run_polling()

if __name__ == "__main__":
    main()
